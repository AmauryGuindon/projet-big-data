"""Génère le support de soutenance Big Data au format PowerPoint (.pptx).

Structure (cahier des charges 5.3) :
- Présentation du problème
- Démonstration technique
- Explication des choix
- Analyse critique du projet
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Emu, Inches, Pt

ROOT = Path(__file__).parent
DIAG = ROOT / "diagrams"
OUTPUT = ROOT / "Soutenance_Big_Data.pptx"

# Palette
COL_PRIMARY = RGBColor(0x2C, 0x3E, 0x50)     # bleu nuit
COL_ACCENT = RGBColor(0xE6, 0x73, 0x22)      # orange
COL_BLUE = RGBColor(0x4A, 0x90, 0xE2)        # bleu vif
COL_GREEN = RGBColor(0x7A, 0xC7, 0x4F)       # vert MinIO
COL_BG_LIGHT = RGBColor(0xF7, 0xF9, 0xFC)    # fond pâle
COL_BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COL_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
COL_MUTED = RGBColor(0x6B, 0x7B, 0x8C)
COL_LINE = RGBColor(0xE3, 0xE8, 0xEF)

# Présentation 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]  # layout vierge


# ============================================================
# Helpers
# ============================================================
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=COL_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0)
    tf.margin_right = Cm(0)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=COL_TEXT, bullet_color=COL_ACCENT,
                line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0)
    tf.margin_top = Cm(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)

        # Puce
        r0 = p.add_run()
        r0.text = "▸ "
        r0.font.name = "Calibri"
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color

        if isinstance(item, tuple):
            label, rest = item
            r1 = p.add_run()
            r1.text = label
            r1.font.name = "Calibri"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = COL_PRIMARY

            r2 = p.add_run()
            r2.text = " — " + rest
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_image_center(slide, filename, top, height_in=None, width_in=None):
    path = str(DIAG / filename)
    if height_in:
        pic = slide.shapes.add_picture(path, 0, top, height=Inches(height_in))
    else:
        pic = slide.shapes.add_picture(path, 0, top, width=Inches(width_in))
    pic.left = int((SW - pic.width) / 2)
    return pic


def slide_header(slide, title, eyebrow=None, page_num=None, total=None):
    """En-tête commun : bande accent, eyebrow + titre."""
    # Bande accent à gauche
    add_rect(slide, 0, 0, Inches(0.18), SH, COL_ACCENT)
    # Bande supérieure subtile
    add_rect(slide, Inches(0.18), 0, SW - Inches(0.18), Inches(0.08), COL_PRIMARY)

    if eyebrow:
        add_text(slide, Inches(0.55), Inches(0.25), Inches(8), Inches(0.4),
                 eyebrow.upper(), size=11, bold=True, color=COL_ACCENT)

    add_text(slide, Inches(0.55), Inches(0.55), Inches(12), Inches(0.9),
             title, size=30, bold=True, color=COL_PRIMARY)

    # Ligne de séparation
    add_rect(slide, Inches(0.55), Inches(1.45), Inches(1.5), Emu(20000), COL_ACCENT)

    # Pagination en bas
    if page_num and total:
        add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                 f"{page_num} / {total}", size=10, color=COL_MUTED,
                 align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.3),
             "Projet Big Data — M1 BDIA — A. Guindon", size=10,
             italic=True, color=COL_MUTED)


def add_notes(slide, text):
    """Ajoute des notes de présentateur sur une diapositive."""
    notes = slide.notes_slide.notes_text_frame
    notes.text = text
    for p in notes.paragraphs:
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(13)


def add_kpi_card(slide, x, y, w, h, value, label, color=COL_BLUE):
    add_rect(slide, x, y, w, h, COL_BG_WHITE, line=COL_LINE)
    # Bande supérieure colorée
    add_rect(slide, x, y, w, Inches(0.12), color)
    add_text(slide, x, y + Inches(0.35), w, Inches(1.0),
             value, size=44, bold=True, color=color, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + h - Inches(0.55), w, Inches(0.5),
             label, size=12, color=COL_MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# Liste des slides à construire (pour la pagination)
# ============================================================
TOTAL_SLIDES = 17


# ============================================================
# Slide 1 — Couverture
# ============================================================
s = prs.slides.add_slide(BLANK)
# Fond dégradé approximatif via deux rectangles
add_rect(s, 0, 0, SW, SH, COL_PRIMARY)
add_rect(s, 0, Inches(5.5), SW, Inches(2), RGBColor(0x1F, 0x2D, 0x3D))

# Accent latéral
add_rect(s, 0, 0, Inches(0.35), SH, COL_ACCENT)

# Eyebrow
add_text(s, Inches(1.0), Inches(1.3), Inches(10), Inches(0.5),
         "PROJET BIG DATA — M1 BIG DATA IA — 2025/2026", size=14, bold=True,
         color=COL_ACCENT)

# Titre principal
add_text(s, Inches(1.0), Inches(2.0), Inches(11.5), Inches(1.5),
         "Plateforme distribuée\nd'ingestion et d'analyse",
         size=54, bold=True, color=COL_BG_WHITE)

# Sous-titre
add_text(s, Inches(1.0), Inches(3.9), Inches(11), Inches(0.6),
         "Cassandra · MinIO (S3) · Docker", size=22,
         color=COL_BLUE, italic=True)

# Auteur en bas
add_text(s, Inches(1.0), Inches(5.9), Inches(11), Inches(0.45),
         "GUINDON Amaury", size=20, bold=True, color=COL_BG_WHITE)
add_text(s, Inches(1.0), Inches(6.35), Inches(11), Inches(0.4),
         "École Supérieure du Numérique 89 — Soutenance orale", size=14,
         color=RGBColor(0xC0, 0xCE, 0xDC))

add_notes(s, (
    "Accroche (30 sec) :\n\n"
    "Bonjour, je suis Amaury Guindon, étudiant en M1 Big Data IA. "
    "Je vais vous présenter mon projet de fin de module : une plateforme "
    "distribuée d'ingestion et d'analyse de données massives, construite "
    "autour de trois briques : Apache Cassandra, un stockage objet S3, "
    "et Docker.\n\n"
    "Objectif de la présentation : démontrer la maîtrise des concepts "
    "Big Data — modélisation NoSQL, scalabilité, théorème CAP — "
    "appliqués à un cas concret et reproductible.\n\n"
    "Durée prévue : environ 15 minutes + démo + questions."
))


# ============================================================
# Slide 2 — Sommaire
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Plan de la présentation", eyebrow="Sommaire",
             page_num=2, total=TOTAL_SLIDES)

cols = [
    ("01", "Contexte & problématique"),
    ("02", "Objectifs"),
    ("03", "Architecture globale"),
    ("04", "Flux de données"),
    ("05", "Choix techniques"),
    ("06", "Modèle Cassandra"),
    ("07", "Principes Big Data — 5V & CAP"),
    ("08", "Démonstration"),
    ("09", "Performances & limites"),
    ("10", "Sécurité & RGPD"),
    ("11", "Pistes d'évolution"),
    ("12", "Conclusion"),
]
# 2 colonnes de 6 lignes
left_x = Inches(0.9)
right_x = Inches(6.9)
y0 = Inches(2.0)
row_h = Inches(0.65)

for idx, (num, title) in enumerate(cols):
    col = idx // 6
    row = idx % 6
    x = left_x if col == 0 else right_x
    y = y0 + row_h * row

    # Numéro
    add_text(s, x, y, Inches(0.9), Inches(0.6),
             num, size=28, bold=True, color=COL_ACCENT,
             anchor=MSO_ANCHOR.MIDDLE)
    # Titre
    add_text(s, x + Inches(0.95), y, Inches(5.0), Inches(0.6),
             title, size=18, color=COL_PRIMARY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

add_notes(s, (
    "Annoncer rapidement le plan (1 min) :\n\n"
    "Je vais structurer la présentation en quatre temps :\n"
    "1) Le contexte et les choix techniques (slides 3 à 7)\n"
    "2) La modélisation Cassandra concrète (slides 8 à 9)\n"
    "3) Les fondements théoriques — 5V et CAP (slides 10 à 11)\n"
    "4) La démo, l'analyse critique et les perspectives (slides 12 à 16)\n\n"
    "Je termine par les questions. Ne pas hésiter à m'interrompre."
))


# ============================================================
# Slide 3 — Contexte & problématique
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Contexte & problématique", eyebrow="01 / Contexte",
             page_num=3, total=TOTAL_SLIDES)

add_text(s, Inches(0.55), Inches(1.7), Inches(12), Inches(0.6),
         "Une entreprise souhaite mettre en place une plateforme capable de :",
         size=18, italic=True, color=COL_MUTED)

add_bullets(s, Inches(0.7), Inches(2.4), Inches(12.0), Inches(4.2), [
    ("Collecter", "des événements hétérogènes — transactions, catalogue produits, contenus multimédia"),
    ("Stocker", "ces données de manière scalable, sans limite haute prévisible"),
    ("Interroger", "rapidement, sur des dimensions métier (pays, marque, catégorie)"),
    ("Garantir", "la disponibilité même en cas de panne d'un nœud"),
    ("Déployer", "de façon reproductible sur n'importe quel environnement"),
], size=18)

add_notes(s, (
    "Poser le problème métier (1 min 30) :\n\n"
    "Imaginer une entreprise qui vend en ligne et qui gère un catalogue mode. "
    "Elle reçoit en continu trois types de données très différents : des "
    "transactions structurées (CSV), un catalogue produits qui évolue (JSON), "
    "et des images de produits (binaires).\n\n"
    "Insister sur les 5 verbes :\n"
    "- Collecter du hétérogène → impose un pipeline multi-format\n"
    "- Stocker scalable → impose une base distribuée, pas un SGBDR\n"
    "- Interroger rapidement → impose un modèle orienté requêtes\n"
    "- Disponibilité → impose la tolérance aux pannes\n"
    "- Reproductibilité → impose la conteneurisation\n\n"
    "Chacun de ces verbes va justifier un choix technique sur les slides "
    "suivantes."
))


# ============================================================
# Slide 4 — Objectifs
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Objectifs du projet", eyebrow="02 / Objectifs",
             page_num=4, total=TOTAL_SLIDES)

# 4 cartes objectifs
card_w, card_h = Inches(2.85), Inches(2.4)
gap = Inches(0.18)
x0 = Inches(0.55)
y_card = Inches(2.0)
cards = [
    ("Architecture", "cohérente, distribuée\net documentée", COL_BLUE),
    ("Modèle Cassandra", "orienté requêtes,\npartition/clustering", COL_ACCENT),
    ("Stockage cloud", "Data Lake S3-compatible\n(MinIO local + AWS)", COL_GREEN),
    ("Déploiement", "Docker reproductible\nen une commande", COL_PRIMARY),
]
for i, (title, desc, color) in enumerate(cards):
    x = x0 + (card_w + gap) * i
    add_rect(s, x, y_card, card_w, card_h, COL_BG_WHITE, line=COL_LINE)
    add_rect(s, x, y_card, card_w, Inches(0.12), color)
    add_text(s, x + Inches(0.2), y_card + Inches(0.35), card_w - Inches(0.4),
             Inches(0.6), title, size=18, bold=True, color=color)
    add_text(s, x + Inches(0.2), y_card + Inches(1.0), card_w - Inches(0.4),
             Inches(1.3), desc, size=14, color=COL_TEXT)

# Pied : analyse Big Data
add_text(s, Inches(0.55), Inches(5.0), Inches(12), Inches(0.5),
         "+ Argumentation Big Data ancrée dans les principes 5V & CAP",
         size=18, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

add_notes(s, (
    "Transition rapide (45 sec) :\n\n"
    "Pour répondre à ces besoins, le projet livre concrètement quatre choses :\n\n"
    "1) Une architecture cohérente — pas juste une base, mais un pipeline complet\n"
    "2) Un modèle Cassandra rigoureux, orienté requêtes (point différenciant vs SQL)\n"
    "3) Un stockage cloud — MinIO en local, AWS S3 en prod, même code\n"
    "4) Un déploiement reproductible — une seule commande Docker\n\n"
    "Et au-delà du code, une argumentation théorique : les 5V et le théorème CAP.\n\n"
    "On va maintenant entrer dans le détail de l'architecture."
))


# ============================================================
# Slide 5 — Architecture globale (avec diagramme)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Architecture globale", eyebrow="03 / Architecture",
             page_num=5, total=TOTAL_SLIDES)

add_image_center(s, "01_architecture.png", top=Inches(1.6), height_in=4.5)
add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
         "Pipeline Python → écriture Cassandra (chaud) + Data Lake MinIO (froid) — orchestré par Docker Compose",
         size=14, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

add_notes(s, (
    "Présenter le schéma de gauche à droite (1 min 30) :\n\n"
    "Tout est dans le cadre Docker — c'est le point clé : un évaluateur lance "
    "'docker compose up' et tout démarre.\n\n"
    "À gauche : les trois sources, les trois formats du périmètre (CSV, JSON, JPG).\n\n"
    "Au centre : le service d'ingestion Python. C'est lui qui lit, valide, "
    "convertit, et distribue vers deux destinations.\n\n"
    "À droite, deux destinations complémentaires :\n"
    "- Cassandra → données 'chaudes' interrogeables en CQL\n"
    "- MinIO (S3) → données 'froides', archives, raw zone\n\n"
    "Pourquoi les deux ? Parce qu'on sépare stockage et calcul — principe "
    "fondamental d'une plateforme Big Data. Si demain on veut rejouer le "
    "pipeline avec un nouveau modèle Cassandra, on n'a pas besoin de "
    "redemander les données aux sources."
))


# ============================================================
# Slide 6 — Flux de données (avec diagramme)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Flux d'ingestion", eyebrow="04 / Flux de données",
             page_num=6, total=TOTAL_SLIDES)

add_image_center(s, "02_flux_donnees.png", top=Inches(1.9), height_in=3.6)

add_bullets(s, Inches(1.0), Inches(5.6), Inches(11.5), Inches(1.5), [
    "Pipeline one-shot, écrit en Python — convient au MVP",
    "Insertions Cassandra par batch de 100 lignes (BatchStatement)",
    "Évolutif vers un consommateur Kafka pour le quasi-temps réel",
], size=15)

add_notes(s, (
    "Détail du pipeline (1 min) :\n\n"
    "Cinq étapes, classiques d'un ETL :\n\n"
    "1) Lecture en streaming — on ne charge pas tout en RAM\n"
    "2) Parsing typé — conversion des dates, des décimaux, gestion des nulls\n"
    "3) Préparation en batch de 100 lignes — amortit le coût réseau\n"
    "4) Insertion Cassandra avec BatchStatement\n"
    "5) En parallèle : upload du fichier brut vers MinIO\n\n"
    "Point important : c'est un job 'one-shot'. Il s'exécute, ingère, "
    "et s'arrête. Pour un MVP c'est suffisant. En production on remplacerait "
    "ça par un consommateur Kafka — sans toucher au modèle Cassandra. "
    "C'est ça la force du modèle orienté requêtes : il découple le pipeline "
    "d'ingestion du stockage."
))


# ============================================================
# Slide 7 — Stack & choix techniques
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Stack & choix techniques", eyebrow="05 / Choix techniques",
             page_num=7, total=TOTAL_SLIDES)

# 4 colonnes : composant + "pourquoi"
col_w = Inches(2.95)
col_gap = Inches(0.15)
x0 = Inches(0.55)
y_top = Inches(2.0)

stack = [
    ("Apache\nCassandra 4.1",
     ["NoSQL distribué",
      "Scalabilité linéaire",
      "Modèle AP (CAP)",
      "Requêtes ms"],
     COL_BLUE),
    ("MinIO\n(S3-compatible)",
     ["Stockage objet",
      "Sépare stockage/calcul",
      "Bascule AWS S3 sans modif",
      "Open source, gratuit"],
     COL_GREEN),
    ("Python 3.11\n(cassandra-driver, boto3)",
     ["Écosystème mature",
      "Drivers officiels",
      "Lisible & maintenable",
      "Job batch idéal pour MVP"],
     COL_ACCENT),
    ("Docker +\nCompose",
     ["Reproductibilité",
      "Healthchecks intégrés",
      "Volumes persistants",
      "Démarrage ordonné"],
     COL_PRIMARY),
]

for i, (title, items, color) in enumerate(stack):
    x = x0 + (col_w + col_gap) * i
    add_rect(s, x, y_top, col_w, Inches(4.5), COL_BG_WHITE, line=COL_LINE)
    add_rect(s, x, y_top, col_w, Inches(0.7), color)
    add_text(s, x + Inches(0.1), y_top + Inches(0.05), col_w - Inches(0.2),
             Inches(0.6), title, size=14, bold=True, color=COL_BG_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Bullets
    tb = s.shapes.add_textbox(x + Inches(0.15), y_top + Inches(0.85),
                              col_w - Inches(0.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = "• "
        r1.font.size = Pt(13)
        r1.font.color.rgb = color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = it
        r2.font.size = Pt(13)
        r2.font.color.rgb = COL_TEXT

add_notes(s, (
    "Parcourir les 4 colonnes (1 min 30) :\n\n"
    "CASSANDRA — pourquoi pas PostgreSQL ?\n"
    "Parce qu'au-delà d'un certain volume, un SGBDR ne scale plus "
    "horizontalement de manière transparente. Cassandra a été pensée pour ça "
    "dès l'origine, par Facebook puis Apache.\n\n"
    "MINIO — pourquoi pas un simple disque ?\n"
    "Parce qu'un Data Lake doit gérer le téraoctet et accepter une API "
    "standard. MinIO parle S3 — donc bascule vers AWS sans changer une ligne "
    "de code applicatif. Juste S3_ENDPOINT_URL à neutraliser.\n\n"
    "PYTHON — pourquoi pas Java ou Scala ?\n"
    "Pour le MVP, le coût d'entrée est minimal et les drivers sont "
    "excellents. En production on pourrait passer à Spark.\n\n"
    "DOCKER — pourquoi pas une install native ?\n"
    "Parce que le cahier des charges impose la reproductibilité. Docker "
    "Compose orchestre l'ordre de démarrage avec les healthchecks — "
    "ingestor ne démarre que quand Cassandra est prête."
))


# ============================================================
# Slide 8 — Modélisation Cassandra (diagramme)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Modélisation orientée requêtes", eyebrow="06 / Modèle Cassandra",
             page_num=8, total=TOTAL_SLIDES)

add_image_center(s, "03_modele_cassandra.png", top=Inches(1.6), height_in=4.2)

add_text(s, Inches(0.55), Inches(6.1), Inches(12.3), Inches(0.5),
         "Une table = une requête   •   PK composite pour répartir   •   CK pour ordonner",
         size=16, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER, bold=True)

add_text(s, Inches(0.55), Inches(6.55), Inches(12.3), Inches(0.4),
         "Keyspace bigdata_project — SimpleStrategy{1} en dev, NetworkTopologyStrategy{3,3} en prod",
         size=13, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

add_notes(s, (
    "Point central — bien expliquer (2 min) :\n\n"
    "La grosse différence avec SQL : en Cassandra on ne part PAS des entités, "
    "on part des REQUÊTES.\n\n"
    "Trois requêtes métier → trois tables :\n\n"
    "TABLE 1 — sales_by_country_month\n"
    "Pour répondre à 'ventes d'un pays sur un mois'. PK = (country, year_month). "
    "Toutes les ventes d'un même pays/mois sont co-localisées sur le même nœud → "
    "une seule partition lue.\n\n"
    "TABLE 2 — products_by_brand_availability\n"
    "PK = (brand, out_of_stock). Cas d'usage : service de réapprovisionnement "
    "qui veut tous les produits en rupture d'une marque.\n\n"
    "TABLE 3 — image_metadata_by_category_gender\n"
    "Le binaire vit dans S3, Cassandra stocke juste la clé S3 + les métadonnées. "
    "Pattern classique : métadonnées chaudes / binaires froids.\n\n"
    "Réplication : SimpleStrategy RF=1 en dev (mono-nœud), "
    "NetworkTopologyStrategy RF=3 en prod multi-DC."
))


# ============================================================
# Slide 9 — Exemple CQL
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Exemple : requête analytique", eyebrow="06 / Modèle Cassandra",
             page_num=9, total=TOTAL_SLIDES)

# Table cible
add_text(s, Inches(0.55), Inches(1.7), Inches(12), Inches(0.5),
         "Requête : ventes du Royaume-Uni en décembre 2010, triées par date",
         size=18, color=COL_MUTED, italic=True)

# Bloc CQL
cql = ("SELECT invoice_date, invoice_no, stock_code,\n"
       "       quantity, unit_price, total_amount\n"
       "FROM sales_by_country_month\n"
       "WHERE country = 'United Kingdom'\n"
       "  AND year_month = '2010-12'\n"
       "LIMIT 20;")

add_rect(s, Inches(0.55), Inches(2.4), Inches(12.2), Inches(2.5),
         RGBColor(0x1E, 0x29, 0x36))
add_text(s, Inches(0.8), Inches(2.55), Inches(11.8), Inches(2.3),
         cql, size=22, color=COL_BG_WHITE, font="Consolas", bold=False)

# Pourquoi c'est rapide
add_bullets(s, Inches(0.55), Inches(5.2), Inches(12), Inches(1.8), [
    ("PK ciblée", "la requête atteint une seule partition → latence ~ ms"),
    ("Tri natif", "CLUSTERING ORDER BY invoice_date DESC, pas de tri applicatif"),
    ("Constant", "le coût ne dépend pas du volume total de la table"),
], size=16)

add_notes(s, (
    "Illustrer concrètement (1 min) :\n\n"
    "On clause WHERE sur les deux composantes de la clé de partition : "
    "country ET year_month. Cassandra exige cela — sans la PK complète, "
    "la requête ne passe pas (sauf à utiliser ALLOW FILTERING, qu'on évite).\n\n"
    "Trois propriétés à souligner :\n\n"
    "1) Une seule partition ciblée → latence quasi constante en millisecondes, "
    "même avec un milliard de lignes en base.\n\n"
    "2) Tri automatique grâce à CLUSTERING ORDER BY invoice_date DESC. "
    "Pas besoin d'ORDER BY applicatif — c'est gratuit.\n\n"
    "3) Cette latence ne dégrade PAS quand le volume total augmente. "
    "C'est la grosse différence avec un SGBDR où la latence croît avec "
    "la taille de la table.\n\n"
    "Si on demande pourquoi pas un index ? → un index secondaire Cassandra "
    "scanne potentiellement tous les nœuds. Mauvaise pratique."
))


# ============================================================
# Slide 10 — Principes Big Data : 5V
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Les 5V appliqués au projet", eyebrow="07 / Big Data",
             page_num=10, total=TOTAL_SLIDES)

five_v = [
    ("VOLUME", "Cassandra + S3\ndimensionnés To+", COL_BLUE),
    ("VÉLOCITÉ", "Batch 100 lignes,\névolutif vers Kafka", COL_ACCENT),
    ("VARIÉTÉ", "CSV · JSON · JPG\n3 pipelines dédiés", COL_GREEN),
    ("VÉRACITÉ", "Parsing défensif\n+ raw zone S3", COL_PRIMARY),
    ("VALEUR", "Tables orientées\nrequêtes métier", RGBColor(0xC0, 0x39, 0x2B)),
]

card_w = Inches(2.4)
gap = Inches(0.12)
x0 = Inches(0.45)
y_card = Inches(2.4)

for i, (title, desc, color) in enumerate(five_v):
    x = x0 + (card_w + gap) * i
    add_rect(s, x, y_card, card_w, Inches(3.2), COL_BG_WHITE, line=COL_LINE)
    add_rect(s, x, y_card, card_w, Inches(1.2), color)
    add_text(s, x, y_card + Inches(0.2), card_w, Inches(0.8),
             title, size=28, bold=True, color=COL_BG_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), y_card + Inches(1.4), card_w - Inches(0.3),
             Inches(1.7), desc, size=14, color=COL_TEXT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_notes(s, (
    "Cadre théorique — rapide (1 min) :\n\n"
    "Les 5V sont le canon pour caractériser un projet Big Data. "
    "Mon projet adresse chacun :\n\n"
    "VOLUME — Cassandra + S3 dimensionnés au téraoctet. "
    "Les limites actuelles (CSV_LIMIT, etc.) sont des garde-fous, pas des plafonds.\n\n"
    "VÉLOCITÉ — Pour l'instant batch, mais batch optimisé (BatchStatement par 100). "
    "Migration vers Kafka possible sans refonte.\n\n"
    "VARIÉTÉ — c'est le V le plus visible : 3 formats, 3 pipelines, 3 tables.\n\n"
    "VÉRACITÉ — parsing défensif côté ingestion + conservation des données "
    "brutes sur S3 → on peut TOUJOURS rejouer.\n\n"
    "VALEUR — les tables ne sont pas neutres, elles répondent à des questions "
    "métier précises. C'est la finalité.\n\n"
    "Transition : 'parmi ces 5V, le plus délicat reste la disponibilité — "
    "c'est ce qui mène au théorème CAP'."
))


# ============================================================
# Slide 11 — CAP Theorem
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Théorème CAP — pourquoi Cassandra ?",
             eyebrow="07 / Big Data", page_num=11, total=TOTAL_SLIDES)

# Image à gauche
pic = s.shapes.add_picture(str(DIAG / "04_cap_theorem.png"),
                            Inches(0.6), Inches(1.7),
                            height=Inches(5.0))

# Texte à droite
add_text(s, Inches(7.0), Inches(1.9), Inches(5.8), Inches(0.6),
         "Cassandra = AP", size=26, bold=True, color=COL_BLUE)
add_bullets(s, Inches(7.0), Inches(2.7), Inches(5.8), Inches(4.0), [
    ("Availability", "le service écrit/lit même en cas de partition réseau"),
    ("Partition tolerance", "natif dans tout système distribué"),
    ("Consistency", "éventuelle, réglable par requête (ONE / QUORUM / ALL)"),
    ("Cas d'usage", "préférable à un downtime pour un service de collecte"),
], size=15)

add_notes(s, (
    "Théorème CAP — point sensible aux questions du jury (1 min 30) :\n\n"
    "Brewer (2000) : dans un système distribué, on ne peut garantir "
    "simultanément que 2 propriétés sur 3 : Consistency, Availability, "
    "Partition tolerance.\n\n"
    "La partition réseau est inévitable dans un système distribué — "
    "donc le vrai choix est entre C et A.\n\n"
    "Cassandra a CHOISI A. Quand un nœud n'arrive plus à parler à un autre, "
    "il continue de servir avec ce qu'il sait. La cohérence se rétablit "
    "ensuite (eventual consistency, via read repair et hinted handoff).\n\n"
    "Important — Cassandra est NUANCÉ : on peut régler la cohérence par "
    "REQUÊTE avec les Consistency Levels :\n"
    "  • ONE → max disponibilité, cohérence faible\n"
    "  • QUORUM → équilibre (majorité de nœuds doivent confirmer)\n"
    "  • ALL → cohérence forte, mais plus fragile\n\n"
    "Pour un service de collecte d'événements, A > C : il vaut mieux une "
    "donnée légèrement décalée qu'un downtime."
))


# ============================================================
# Slide 12 — Démonstration technique
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Démonstration technique", eyebrow="08 / Démo",
             page_num=12, total=TOTAL_SLIDES)

# Étapes de démo
add_text(s, Inches(0.55), Inches(1.7), Inches(12), Inches(0.5),
         "Déroulé live (≈ 3 min) :",
         size=18, italic=True, color=COL_MUTED)

steps = [
    ("1. Lancer la stack",
     "docker compose up --build"),
    ("2. Suivre l'ingestion",
     "logs ingestor : raw uploads → batch Cassandra"),
    ("3. Vérifier Cassandra",
     "docker exec -it bd_cassandra cqlsh + SELECT COUNT(*)"),
    ("4. Exécuter une requête métier",
     "SELECT … WHERE country='United Kingdom' AND year_month='2010-12'"),
    ("5. Vérifier le Data Lake",
     "console MinIO http://localhost:9001 — bucket bigdata-raw"),
]
y = Inches(2.4)
for title, cmd in steps:
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.8),
             COL_BG_WHITE, line=COL_LINE)
    add_text(s, Inches(0.75), y + Inches(0.05), Inches(4.5), Inches(0.7),
             title, size=15, bold=True, color=COL_ACCENT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.3), y + Inches(0.05), Inches(7.3), Inches(0.7),
             cmd, size=13, color=COL_TEXT, font="Consolas",
             anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.88)

add_notes(s, (
    "Démo live (3 min — PRÉPARER À L'AVANCE) :\n\n"
    "AVANT LA SOUTENANCE :\n"
    "- Faire un 'docker compose down -v' pour repartir propre\n"
    "- Avoir un terminal PowerShell ouvert dans le dossier projet\n"
    "- Avoir le navigateur prêt sur http://localhost:9001\n\n"
    "DÉROULÉ :\n\n"
    "1) Lancer 'docker compose up --build' (peut prendre ~1 min au premier "
    "lancement à cause du pull Cassandra) — pendant ce temps, on peut "
    "continuer à expliquer.\n\n"
    "2) Montrer les logs : on voit Cassandra ready → init schema → "
    "ingestor démarre → uploads S3 puis batches Cassandra.\n\n"
    "3) Ouvrir cqlsh dans Cassandra et lancer un COUNT(*) pour montrer "
    "que les données sont bien là.\n\n"
    "4) Lancer la requête métier (slide 9) en direct → résultat instantané.\n\n"
    "5) Bascule sur la console MinIO → montrer le bucket bigdata-raw "
    "avec les 3 zones (structured / semi_structured / non_structured) "
    "+ les images.\n\n"
    "PLAN B : si la démo coince, screenshots de secours dans le dossier "
    "rapport/."
))


# ============================================================
# Slide 13 — Performances : KPI
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Ce que mesure le pipeline", eyebrow="09 / Performances",
             page_num=13, total=TOTAL_SLIDES)

# 4 KPI
add_kpi_card(s, Inches(0.6), Inches(2.0), Inches(2.9), Inches(2.4),
             "3", "familles de données\n(structuré / semi / non)", COL_BLUE)
add_kpi_card(s, Inches(3.65), Inches(2.0), Inches(2.9), Inches(2.4),
             "3", "tables Cassandra\norientées requêtes", COL_ACCENT)
add_kpi_card(s, Inches(6.7), Inches(2.0), Inches(2.9), Inches(2.4),
             "100", "lignes / batch\n(BatchStatement)", COL_GREEN)
add_kpi_card(s, Inches(9.75), Inches(2.0), Inches(2.9), Inches(2.4),
             "1", "commande pour\ntout déployer", COL_PRIMARY)

# Bandeau "points forts"
add_text(s, Inches(0.55), Inches(4.8), Inches(12), Inches(0.5),
         "Points forts", size=20, bold=True, color=COL_PRIMARY)
add_bullets(s, Inches(0.7), Inches(5.3), Inches(12), Inches(2.0), [
    "Latence ms — toute requête principale tape une seule partition",
    "Débit d'écriture élevé — batchs préparés (PreparedStatement)",
    "Découplage stockage / calcul — retraitement complet depuis la raw zone",
    "Stack 100 % open source — coût nul, reproductibilité totale",
], size=15)

add_notes(s, (
    "Récapitulatif quantitatif (45 sec) :\n\n"
    "Quelques chiffres clés pour résumer ce qui est livré :\n\n"
    "- 3 familles de données couvertes (structuré / semi / non)\n"
    "- 3 tables Cassandra, chacune dédiée à une requête métier\n"
    "- 100 lignes par BatchStatement (compromis perf / mémoire)\n"
    "- 1 seule commande pour tout démarrer\n\n"
    "Côté qualitatif, 4 points forts à retenir :\n\n"
    "- Latence en millisecondes — toute requête principale = 1 partition\n"
    "- Débit d'écriture élevé grâce aux PreparedStatement (la requête CQL "
    "est compilée une seule fois côté Cassandra)\n"
    "- Découplage stockage/calcul — on peut détruire Cassandra et la "
    "reconstruire depuis S3\n"
    "- 100% open source : zéro coût de licence"
))


# ============================================================
# Slide 14 — Limites
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Limites & analyse critique", eyebrow="09 / Analyse critique",
             page_num=14, total=TOTAL_SLIDES)

# 2 colonnes : ce qui marche / ce qui manque
left_x = Inches(0.55)
right_x = Inches(6.85)
col_w = Inches(6.0)
y_top = Inches(1.9)

add_rect(s, left_x, y_top, col_w, Inches(5.0), COL_BG_WHITE, line=COL_LINE)
add_rect(s, left_x, y_top, col_w, Inches(0.6), COL_GREEN)
add_text(s, left_x, y_top, col_w, Inches(0.6),
         "✓  Ce qui est démontré", size=18, bold=True, color=COL_BG_WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, left_x + Inches(0.2), y_top + Inches(0.8),
            col_w - Inches(0.4), Inches(4.0), [
    "Modélisation orientée requêtes",
    "Réplication paramétrable",
    "Conteneurisation reproductible",
    "Trois familles de données ingérées",
    "Raw zone + couche chaude séparées",
], size=14)

add_rect(s, right_x, y_top, col_w, Inches(5.0), COL_BG_WHITE, line=COL_LINE)
add_rect(s, right_x, y_top, col_w, Inches(0.6), RGBColor(0xC0, 0x39, 0x2B))
add_text(s, right_x, y_top, col_w, Inches(0.6),
         "△  Ce qui reste à faire", size=18, bold=True, color=COL_BG_WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, right_x + Inches(0.2), y_top + Inches(0.8),
            col_w - Inches(0.4), Inches(4.0), [
    "Mono-nœud Cassandra (RF=1)",
    "Pas de streaming temps réel",
    "Pas de couche analytique (Spark)",
    "Batch parfois multi-partitions",
    "Sécurité non activée (dev)",
], size=14, bullet_color=RGBColor(0xC0, 0x39, 0x2B))

add_notes(s, (
    "Honnêteté technique (1 min 30) :\n\n"
    "Le jury va apprécier de voir qu'on connaît les LIMITES de ce qu'on livre.\n\n"
    "À GAUCHE — ce qui est démontré :\n"
    "Toute la chaîne fonctionne, la modélisation est rigoureuse, la "
    "conteneurisation est propre.\n\n"
    "À DROITE — ce qui manque honnêtement :\n\n"
    "1) MONO-NŒUD : avec un seul nœud Cassandra, je ne peux pas DÉMONTRER "
    "concrètement la tolérance aux pannes. Il faudrait un cluster à 3 nœuds.\n\n"
    "2) BATCH ONLY : pas de Kafka, pas de streaming.\n\n"
    "3) PAS DE SPARK : pour des agrégations cross-partitions (ex: total des "
    "ventes mondiales par produit), il faut un moteur dédié.\n\n"
    "4) BATCHS MULTI-PARTITIONS : techniquement, mes BatchStatement peuvent "
    "toucher plusieurs partitions en un envoi → mauvaise pratique Cassandra "
    "(pression sur le coordinateur). Correction prévue.\n\n"
    "5) SÉCURITÉ : auth Cassandra désactivée, identifiants MinIO par défaut. "
    "OK pour un dev, à durcir en prod (slide suivante)."
))


# ============================================================
# Slide 15 — Sécurité & RGPD
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Sécurité & RGPD", eyebrow="10 / Enjeux transverses",
             page_num=15, total=TOTAL_SLIDES)

add_text(s, Inches(0.55), Inches(1.85), Inches(6), Inches(0.5),
         "Sécurité — à activer en prod", size=20, bold=True, color=COL_BLUE)
add_bullets(s, Inches(0.7), Inches(2.45), Inches(6.0), Inches(4.5), [
    "Authentification Cassandra (PasswordAuthenticator)",
    "TLS client → cluster + inter-nœuds",
    "Rotation des credentials MinIO",
    "Isolation réseau Docker",
    "Audit logs centralisés",
], size=15)

add_text(s, Inches(6.85), Inches(1.85), Inches(6.3), Inches(0.5),
         "RGPD — principes appliqués", size=20, bold=True, color=COL_ACCENT)
add_bullets(s, Inches(7.0), Inches(2.45), Inches(6.0), Inches(4.5), [
    ("Minimisation", "ne stocker que les champs nécessaires"),
    ("Effacement", "TTL plutôt que DELETE (tombstones)"),
    ("Localisation", "NetworkTopologyStrategy = DC EU/US"),
    ("Traçabilité", "audit logs Cassandra + S3"),
    ("Chiffrement", "SSE sur le bucket S3"),
], size=15, bullet_color=COL_ACCENT)

add_notes(s, (
    "Sécurité + RGPD (1 min 30) :\n\n"
    "Pour la SÉCURITÉ, en production :\n"
    "- Authentification Cassandra avec PasswordAuthenticator + "
    "CassandraAuthorizer, compte applicatif aux droits restreints\n"
    "- TLS partout : client→cluster ET inter-nœuds\n"
    "- Rotation des credentials MinIO via Docker secrets ou un vault\n"
    "- Réseau Docker dédié, ports non exposés à l'extérieur\n\n"
    "Pour le RGPD, point sensible avec des données client :\n\n"
    "MINIMISATION — ne pas stocker plus que nécessaire. customer_id en texte "
    "libre aujourd'hui → pseudonymisé en prod.\n\n"
    "DROIT À L'EFFACEMENT — point délicat dans Cassandra : un DELETE crée "
    "un 'tombstone' qui ralentit les lectures. Meilleure pratique : TTL "
    "(Time To Live) sur les lignes contenant des PII, ou séparer dans une "
    "table dédiée plus facile à purger.\n\n"
    "LOCALISATION — NetworkTopologyStrategy permet de garantir que les "
    "données UE restent sur des nœuds européens.\n\n"
    "TRAÇABILITÉ — audit logs des deux côtés (Cassandra + S3).\n\n"
    "CHIFFREMENT — SSE-S3 ou SSE-KMS côté bucket."
))


# ============================================================
# Slide 16 — Pistes d'évolution
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Pistes d'évolution", eyebrow="11 / Évolutions",
             page_num=16, total=TOTAL_SLIDES)

evolutions = [
    ("Cluster 3 nœuds",
     "NetworkTopologyStrategy, RF=3, démonstration concrète de la tolérance aux pannes",
     COL_BLUE),
    ("Streaming Kafka",
     "Remplacer le job batch par un consommateur temps réel — sans modifier le modèle Cassandra",
     COL_ACCENT),
    ("Couche Spark / Trino",
     "Agrégations cross-partitions, jointures, exports vers reporting",
     COL_GREEN),
    ("Observabilité",
     "Prometheus + Grafana — latence, débit, état du cluster",
     COL_PRIMARY),
]

y = Inches(2.0)
for title, desc, color in evolutions:
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(1.05),
             COL_BG_WHITE, line=COL_LINE)
    add_rect(s, Inches(0.55), y, Inches(0.18), Inches(1.05), color)
    add_text(s, Inches(0.95), y + Inches(0.12), Inches(3.5), Inches(0.5),
             title, size=18, bold=True, color=color)
    add_text(s, Inches(0.95), y + Inches(0.55), Inches(11.5), Inches(0.5),
             desc, size=13, color=COL_TEXT)
    y += Inches(1.15)

add_notes(s, (
    "Perspectives (1 min) :\n\n"
    "Quatre axes d'évolution naturels, par ordre de priorité :\n\n"
    "1) CLUSTER 3 NŒUDS — c'est la première chose à faire. RF=3 avec "
    "NetworkTopologyStrategy, et là on démontre VRAIMENT la tolérance aux "
    "pannes : on peut tuer un nœud sans interruption de service.\n\n"
    "2) STREAMING KAFKA — pour les cas d'usage temps réel (clics, IoT). "
    "On ajoute un broker entre les sources et l'ingestor, qui devient un "
    "consommateur Kafka. Le modèle Cassandra ne change pas.\n\n"
    "3) SPARK / TRINO — dès qu'on a besoin de croiser plusieurs tables ou "
    "d'agréger sur l'ensemble du dataset. Spark peut lire directement "
    "depuis Cassandra ET depuis S3.\n\n"
    "4) OBSERVABILITÉ — Prometheus pour les métriques, Grafana pour la "
    "visualisation. Indispensable en prod pour suivre la latence p99, "
    "le débit, l'état du cluster.\n\n"
    "Bonus si questions : Airflow pour orchestrer des pipelines périodiques, "
    "dbt sur Trino pour les transformations, etc."
))


# ============================================================
# Slide 17 — Conclusion / Merci
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, COL_PRIMARY)
add_rect(s, 0, 0, Inches(0.35), SH, COL_ACCENT)

add_text(s, Inches(1.0), Inches(1.6), Inches(11), Inches(0.6),
         "CONCLUSION", size=16, bold=True, color=COL_ACCENT)

add_text(s, Inches(1.0), Inches(2.2), Inches(11), Inches(1.5),
         "Une plateforme Big Data\nfonctionnelle et reproductible.", size=48,
         bold=True, color=COL_BG_WHITE)

add_bullets(s, Inches(1.2), Inches(4.4), Inches(11), Inches(2.0), [
    "Ingestion · stockage distribué · Data Lake · conteneurisation",
    "Choix justifiés par les principes structurants du Big Data",
    "Socle réutilisable, évolutif vers la production",
], size=18, color=COL_BG_WHITE, bullet_color=COL_ACCENT)

add_text(s, Inches(1.0), Inches(6.7), Inches(11), Inches(0.5),
         "Merci pour votre attention — Questions ?", size=20,
         italic=True, color=COL_BLUE, align=PP_ALIGN.CENTER)

add_notes(s, (
    "Conclusion (45 sec) — message à retenir par le jury :\n\n"
    "Le projet livre concrètement les 4 piliers demandés :\n"
    "ingestion, stockage distribué, Data Lake, conteneurisation.\n\n"
    "Chaque choix technique est étayé par les concepts du Big Data — "
    "5V, théorème CAP, scalabilité horizontale — et pas seulement "
    "par 'parce que c'est moderne'.\n\n"
    "Au-delà du périmètre demandé, c'est un SOCLE évolutif vers la "
    "production : ajout d'un cluster multi-DC, passage au streaming, "
    "couche analytique. La fondation tient.\n\n"
    "PRÉPARER LES QUESTIONS PROBABLES :\n\n"
    "Q: Pourquoi pas MongoDB ? → CP, latence moins prédictible sur "
    "des requêtes multi-partitions, modèle document plus souple mais "
    "moins adapté à des requêtes analytiques avec ORDER BY natif.\n\n"
    "Q: Pourquoi pas Elasticsearch ? → orienté recherche full-text, "
    "pas un stockage primaire. Bon en complément, pas en remplacement.\n\n"
    "Q: Si je veux faire un GROUP BY sur toute la base ? → pas Cassandra, "
    "utiliser Spark qui lit la base.\n\n"
    "Q: Pourquoi 100 lignes par batch ? → équilibre empirique. Trop petit "
    "= overhead réseau ; trop grand = risque timeout coordinateur Cassandra "
    "(limite par défaut ~50 ko)."
))


# ============================================================
# Sauvegarde
# ============================================================
prs.save(OUTPUT)
print(f"PowerPoint généré : {OUTPUT}")
